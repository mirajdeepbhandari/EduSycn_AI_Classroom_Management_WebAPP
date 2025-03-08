from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
import os

def generate_ppt(theme_name, output,file_name):
        
        file_name = file_name.replace(".pdf", "")

        introsection = output['FinalResult']['Introduction_']

        bodysection = output['FinalResult']['Body_']

        concsection = output['FinalResult']['Conclusion_']

        if theme_name == 'theme1':
            theme_path = 'static/themes/theme1.pptx'

        elif theme_name == 'theme2':
            theme_path = 'static/themes/theme2.pptx'

        elif theme_name == 'theme3':
            theme_path = 'static/themes/theme3.pptx'
        
        # Load existing presentation
        prs = Presentation(theme_path)

        # Function to add text with separate checks for paragraphs & bullet points
        def add_paginated_text(slide_title, heading_desc, bullet_points):
            """
            Adds text content dynamically, splitting into new slides if necessary.
            - If heading paragraph > 250 words → New slide.
            - If bullet points > 3 → New slide.
            """
            slide_layout = prs.slide_layouts[1]  # Title + Content Layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_title

            content = slide.shapes.placeholders[1]
            text_frame = content.text_frame
            text_frame.word_wrap = True

            # Handling heading description as a paragraph (not a bullet point)
            if heading_desc:
                word_count = len(heading_desc.split())

                if word_count > 100:
                    # Split the text into chunks of approximately 250 words
                    words = heading_desc.split()
                    chunks = []
                    current_chunk = []
                    current_count = 0

                    for word in words:
                        current_chunk.append(word)
                        current_count += 1
                        if current_count >= 100:
                            chunks.append(' '.join(current_chunk))
                            current_chunk = []
                            current_count = 0

                    # Add any remaining words as the last chunk
                    if current_chunk:
                        chunks.append(' '.join(current_chunk))

                    # Add first chunk to current slide
                    p = text_frame.add_paragraph()
                    p.text = chunks[0]
                    p.font.size = Pt(16)
                    p.alignment = PP_ALIGN.JUSTIFY
                    p.level = 0  # Ensure it's not a bullet point
                    p.space_after = Pt(15)

                    # Create new slides for remaining chunks
                    for i in range(1, len(chunks)):
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        title.text = slide_title + " (Continued)"
                        content = slide.shapes.placeholders[1]
                        text_frame = content.text_frame
                        text_frame.word_wrap = True

                        p = text_frame.add_paragraph()
                        p.text = chunks[i]
                        p.font.size = Pt(16)
                        p.alignment = PP_ALIGN.JUSTIFY
                        p.level = 0
                        p.space_after = Pt(15)
                else:
                    # If less than 250 words, just add to current slide
                    p = text_frame.add_paragraph()
                    p.text = heading_desc
                    p.font.size = Pt(16)
                    p.alignment = PP_ALIGN.JUSTIFY
                    p.level = 0  # Ensure it's not a bullet point
                    p.space_after = Pt(15)

            # Handling bullet points separately
            bullet_count = 0  # Track number of bullets
            for bullet in bullet_points:
                if bullet_count >= 4:
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    title.text = slide_title + " (Continued)"
                    content = slide.shapes.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.word_wrap = True
                    bullet_count = 0  # Reset bullet count

                # Add bullet point
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.JUSTIFY
                p.level = 1  # This makes it a bullet point
                p.space_after = Pt(10)

                bullet_count += 1  # Increment bullet count

        # _______________________________________________________________________________________
        # Add a title slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text_frame.text = introsection["MainTitle"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY

        subtitle.text_frame.text = introsection["SubTitle"]
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY

        # _______________________________________________________________________________________
        # Add Introduction Slide
        add_paginated_text("Introduction", introsection["Introduction"], [])

        # _______________________________________________________________________________________
        # For Aims

        add_paginated_text("Aims", "", [introsection['Aims_Objectives'][0]['Aims']])

        # _______________________________________________________________________________________
        # For Objectives
        add_paginated_text("Objectives", "", introsection['Aims_Objectives'][0]['Objectives'])

        # _______________________________________________________________________________________
        # Add body content dynamically from introsection['Pages']
        for page in bodysection['Pages']:
            add_paginated_text(page['Headings'], page.get('Description', ''), page['Bulletpoints'])


        #_______________________________________________________________________________________
        #Add Conclusion
        add_paginated_text("Conclusion", concsection['FinalSummary']['Summary'], [])

        os.makedirs('static/pptStore', exist_ok=True) 
        # Save the presentation
        prs.save(f'static/pptStore/{file_name}.pptx')

        

